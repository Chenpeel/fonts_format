#!/usr/bin/env python3
"""
Office 文件字体标准化工具
- 中文 → 宋体
- 英文 / 数字 → Times New Roman
- 日文假名 → MS Mincho

用法:
  uv run main.py file1.pptx file2.docx file3.xlsx ...
"""

import sys
import copy
import shutil
import subprocess
from pathlib import Path

FONT_CHINESE  = '宋体'
FONT_JAPANESE = 'MS Mincho'
FONT_LATIN    = 'Times New Roman'

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

def qa(tag): return f'{{{NS_A}}}{tag}'
def qw(tag): return f'{{{NS_W}}}{tag}'


# ── 字符分类 ──────────────────────────────────────────────────────────────────

def char_type(ch: str) -> str:
    cp = ord(ch)
    # 日文假名（平假名 / 片假名）
    if 0x3040 <= cp <= 0x309F or 0x30A0 <= cp <= 0x30FF or 0x31F0 <= cp <= 0x31FF:
        return 'japanese'
    # CJK 统一汉字（视为中文）
    if (0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF or
            0x20000 <= cp <= 0x2A6DF or 0xF900 <= cp <= 0xFAFF or
            0x2E80 <= cp <= 0x2EFF):
        return 'chinese'
    # 基本拉丁 + 西欧拉丁（含数字、ASCII 标点）
    if cp < 0x0300:
        return 'latin'
    return 'other'


def segment(text: str) -> list[tuple[str, str]]:
    """将文本按字符类型切分为 [(type, text), ...] 列表"""
    if not text:
        return []
    segs: list[tuple[str, str]] = []
    cur_t, cur_s = char_type(text[0]), text[0]
    for ch in text[1:]:
        t = char_type(ch)
        if t == cur_t:
            cur_s += ch
        else:
            segs.append((cur_t, cur_s))
            cur_t, cur_s = t, ch
    segs.append((cur_t, cur_s))
    return segs


def fonts_for_type(seg_type: str) -> tuple[str, str]:
    """返回 (latin字体, eastAsia字体)，两者保持一致避免 WPS 渲染歧义"""
    if seg_type == 'japanese':
        return FONT_JAPANESE, FONT_JAPANESE
    if seg_type == 'chinese':
        return FONT_CHINESE, FONT_CHINESE
    return FONT_LATIN, FONT_LATIN


# ── PPTX ──────────────────────────────────────────────────────────────────────

def pptx_apply_fonts(r_elem, latin: str, ea: str) -> None:
    from lxml import etree
    rPr = r_elem.find(qa('rPr'))
    if rPr is None:
        rPr = etree.Element(qa('rPr'))
        r_elem.insert(0, rPr)
    for font, tag in [(latin, 'latin'), (ea, 'ea')]:
        el = rPr.find(qa(tag))
        if el is None:
            el = etree.SubElement(rPr, qa(tag))
        el.set('typeface', font)


def pptx_process_para(p_elem) -> None:
    for r_elem in list(p_elem.findall(qa('r'))):
        t_elem = r_elem.find(qa('t'))
        if t_elem is None or not t_elem.text:
            continue

        segs = segment(t_elem.text)

        # 纯单一类型：直接设字体，无需拆分
        if len(segs) == 1:
            latin, ea = fonts_for_type(segs[0][0])
            pptx_apply_fonts(r_elem, latin, ea)
            continue

        # 混合类型 → 拆分 run，每段独立赋字体
        parent = r_elem.getparent()
        idx = list(parent).index(r_elem)
        parent.remove(r_elem)
        for i, (seg_type, seg_text) in enumerate(segs):
            nr = copy.deepcopy(r_elem)
            nr.find(qa('t')).text = seg_text
            latin, ea = fonts_for_type(seg_type)
            pptx_apply_fonts(nr, latin, ea)
            parent.insert(idx + i, nr)


def process_pptx(path: Path) -> None:
    from pptx import Presentation
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    pptx_process_para(para._p)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            pptx_process_para(para._p)
    out = path.parent / f'{path.stem}_font_fixed{path.suffix}'
    prs.save(out)
    print(f'✓ {out}')


# ── DOCX ──────────────────────────────────────────────────────────────────────

def docx_apply_fonts(r_elem, ascii_f: str, ea_f: str) -> None:
    from lxml import etree
    rPr = r_elem.find(qw('rPr'))
    if rPr is None:
        rPr = etree.Element(qw('rPr'))
        r_elem.insert(0, rPr)
    rFonts = rPr.find(qw('rFonts'))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qw('rFonts'))
    rFonts.set(qw('ascii'),    ascii_f)
    rFonts.set(qw('hAnsi'),    ascii_f)
    rFonts.set(qw('eastAsia'), ea_f)


def docx_process_para(p_elem) -> None:
    XML_SPACE = '{http://www.w3.org/XML/1998/namespace}space'
    for r_elem in list(p_elem.findall('.//' + qw('r'))):
        t_elem = r_elem.find(qw('t'))
        if t_elem is None or not t_elem.text:
            continue

        segs = segment(t_elem.text)

        if len(segs) == 1:
            latin, ea = fonts_for_type(segs[0][0])
            docx_apply_fonts(r_elem, latin, ea)
            continue

        parent = r_elem.getparent()
        idx = list(parent).index(r_elem)
        parent.remove(r_elem)
        for i, (seg_type, seg_text) in enumerate(segs):
            nr = copy.deepcopy(r_elem)
            t = nr.find(qw('t'))
            t.text = seg_text
            if seg_text != seg_text.strip():
                t.set(XML_SPACE, 'preserve')
            latin, ea = fonts_for_type(seg_type)
            docx_apply_fonts(nr, latin, ea)
            parent.insert(idx + i, nr)


def process_docx(path: Path) -> None:
    from docx import Document
    doc = Document(path)
    for para in doc.paragraphs:
        docx_process_para(para._p)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    docx_process_para(para._p)
    out = path.parent / f'{path.stem}_font_fixed{path.suffix}'
    doc.save(out)
    print(f'✓ {out}')


# ── XLSX ──────────────────────────────────────────────────────────────────────

def process_xlsx(path: Path) -> None:
    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.load_workbook(path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if not cell.value or not isinstance(cell.value, str):
                    continue
                types = [char_type(c) for c in cell.value]
                cnt = {'chinese': 0, 'japanese': 0, 'latin': 0, 'other': 0}
                for t in types:
                    cnt[t] += 1
                dominant = max(cnt, key=lambda k: cnt[k])

                sz   = (cell.font.size if cell.font and cell.font.size else 11)
                bold = (cell.font.bold if cell.font else False)

                name = (FONT_JAPANESE if dominant == 'japanese'
                        else FONT_CHINESE if dominant == 'chinese'
                        else FONT_LATIN)
                cell.font = Font(name=name, size=sz, bold=bold)

    out = path.parent / f'{path.stem}_font_fixed{path.suffix}'
    wb.save(out)
    print(f'✓ {out}')


# ── 字体安装 ──────────────────────────────────────────────────────────────────

def ensure_fonts_installed() -> None:
    """将 ./fonts/ 目录中的字体安装到用户字体目录（仅复制缺失的）"""
    script_dir = Path(__file__).parent
    src_dir = script_dir / 'fonts'
    if not src_dir.exists():
        return

    dst_dir = Path.home() / '.local' / 'share' / 'fonts' / 'office-fix'
    dst_dir.mkdir(parents=True, exist_ok=True)

    installed = []
    for src in src_dir.iterdir():
        if src.suffix.lower() not in ('.ttf', '.ttc', '.otf'):
            continue
        dst = dst_dir / src.name
        if not dst.exists():
            shutil.copy2(src, dst)
            installed.append(src.name)

    if installed:
        subprocess.run(['fc-cache', '-f', str(dst_dir)],
                       capture_output=True, check=False)
        print(f'已安装字体: {", ".join(installed)}')


# ── 入口 ──────────────────────────────────────────────────────────────────────

HANDLERS = {
    '.pptx': process_pptx,
    '.docx': process_docx,
    '.xlsx': process_xlsx,
}


def main() -> None:
    ensure_fonts_installed()

    if len(sys.argv) < 2:
        print('用法: uv run main.py file1.pptx [file2.docx file3.xlsx ...]')
        sys.exit(1)

    for arg in sys.argv[1:]:
        p = Path(arg)
        if not p.exists():
            print(f'✗ 文件不存在: {p}')
            continue
        handler = HANDLERS.get(p.suffix.lower())
        if not handler:
            print(f'✗ 不支持的格式: {p.suffix}（支持 .pptx .docx .xlsx）')
            continue
        try:
            print(f'处理: {p.name}')
            handler(p)
        except Exception as e:
            print(f'✗ 处理失败 {p.name}: {e}')
            import traceback
            traceback.print_exc()


if __name__ == '__main__':
    main()
