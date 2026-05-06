#!/usr/bin/env python3
"""字体处理核心逻辑（与界面无关）"""

import copy
import shutil
import subprocess
from pathlib import Path
from typing import Callable

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

def qa(tag): return f'{{{NS_A}}}{tag}'
def qw(tag): return f'{{{NS_W}}}{tag}'

DEFAULT_FONTS = {
    'chinese':  '宋体',
    'japanese': 'MS Mincho',
    'latin':    'Times New Roman',
}


# ── 字符分类 ──────────────────────────────────────────────────────────────────

def char_type(ch: str) -> str:
    cp = ord(ch)
    if 0x3040 <= cp <= 0x309F or 0x30A0 <= cp <= 0x30FF or 0x31F0 <= cp <= 0x31FF:
        return 'japanese'
    if (0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF or
            0x20000 <= cp <= 0x2A6DF or 0xF900 <= cp <= 0xFAFF or
            0x2E80 <= cp <= 0x2EFF):
        return 'chinese'
    if cp < 0x0300:
        return 'latin'
    return 'other'


def segment(text: str) -> list[tuple[str, str]]:
    if not text:
        return []
    segs: list[tuple[str, str]] = []
    cur_t, cur_s = char_type(text[0]), text[0]
    for ch in text[1:]:
        t = char_type(ch)
        if t == cur_t:
            cur_s += ch
        else:
            segs.append((cur_t, cur_s))
            cur_t, cur_s = t, ch
    segs.append((cur_t, cur_s))
    return segs


def resolve_font(seg_type: str, fonts: dict) -> tuple[str, str]:
    """同时设 latin 和 ea，避免 WPS Linux 只读 latin 导致全用同一字体"""
    f = fonts.get(seg_type) or fonts.get('latin', 'Times New Roman')
    return f, f


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _pptx_set(r_elem, latin: str, ea: str) -> None:
    from lxml import etree
    rPr = r_elem.find(qa('rPr'))
    if rPr is None:
        rPr = etree.Element(qa('rPr'))
        r_elem.insert(0, rPr)
    for font, tag in [(latin, 'latin'), (ea, 'ea')]:
        el = rPr.find(qa(tag))
        if el is None:
            el = etree.SubElement(rPr, qa(tag))
        el.set('typeface', font)


def _pptx_para(p_elem, fonts: dict) -> None:
    for r_elem in list(p_elem.findall(qa('r'))):
        t_elem = r_elem.find(qa('t'))
        if t_elem is None or not t_elem.text:
            continue
        segs = segment(t_elem.text)
        if len(segs) == 1:
            _pptx_set(r_elem, *resolve_font(segs[0][0], fonts))
            continue
        parent = r_elem.getparent()
        idx = list(parent).index(r_elem)
        parent.remove(r_elem)
        for i, (st, txt) in enumerate(segs):
            nr = copy.deepcopy(r_elem)
            nr.find(qa('t')).text = txt
            _pptx_set(nr, *resolve_font(st, fonts))
            parent.insert(idx + i, nr)


def process_pptx(path: Path, out_dir: Path, fonts: dict, log: Callable) -> Path:
    from pptx import Presentation
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _pptx_para(para._p, fonts)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _pptx_para(para._p, fonts)
    out = out_dir / f'{path.stem}_font_fixed{path.suffix}'
    prs.save(out)
    return out


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _docx_set(r_elem, ascii_f: str, ea_f: str) -> None:
    from lxml import etree
    rPr = r_elem.find(qw('rPr'))
    if rPr is None:
        rPr = etree.Element(qw('rPr'))
        r_elem.insert(0, rPr)
    rFonts = rPr.find(qw('rFonts'))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qw('rFonts'))
    rFonts.set(qw('ascii'),    ascii_f)
    rFonts.set(qw('hAnsi'),   ascii_f)
    rFonts.set(qw('eastAsia'), ea_f)


def _docx_para(p_elem, fonts: dict) -> None:
    XML_SPACE = '{http://www.w3.org/XML/1998/namespace}space'
    for r_elem in list(p_elem.findall('.//' + qw('r'))):
        t_elem = r_elem.find(qw('t'))
        if t_elem is None or not t_elem.text:
            continue
        segs = segment(t_elem.text)
        if len(segs) == 1:
            _docx_set(r_elem, *resolve_font(segs[0][0], fonts))
            continue
        parent = r_elem.getparent()
        idx = list(parent).index(r_elem)
        parent.remove(r_elem)
        for i, (st, txt) in enumerate(segs):
            nr = copy.deepcopy(r_elem)
            t = nr.find(qw('t'))
            t.text = txt
            if txt != txt.strip():
                t.set(XML_SPACE, 'preserve')
            _docx_set(nr, *resolve_font(st, fonts))
            parent.insert(idx + i, nr)


def process_docx(path: Path, out_dir: Path, fonts: dict, log: Callable) -> Path:
    from docx import Document
    doc = Document(path)
    for para in doc.paragraphs:
        _docx_para(para._p, fonts)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _docx_para(para._p, fonts)
    out = out_dir / f'{path.stem}_font_fixed{path.suffix}'
    doc.save(out)
    return out


# ── XLSX ──────────────────────────────────────────────────────────────────────

def process_xlsx(path: Path, out_dir: Path, fonts: dict, log: Callable) -> Path:
    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.load_workbook(path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if not cell.value or not isinstance(cell.value, str):
                    continue
                cnt = {'chinese': 0, 'japanese': 0, 'latin': 0, 'other': 0}
                for ch in cell.value:
                    cnt[char_type(ch)] += 1
                dominant = max(cnt, key=lambda k: cnt[k])
                sz   = cell.font.size if cell.font and cell.font.size else 11
                bold = cell.font.bold if cell.font else False
                name = fonts.get(dominant) or fonts.get('latin', 'Times New Roman')
                cell.font = Font(name=name, size=sz, bold=bold)

    out = out_dir / f'{path.stem}_font_fixed{path.suffix}'
    wb.save(out)
    return out


# ── 统一入口 ──────────────────────────────────────────────────────────────────

HANDLERS: dict[str, Callable] = {
    '.pptx': process_pptx,
    '.docx': process_docx,
    '.xlsx': process_xlsx,
}


def process_file(path: Path, out_dir: Path | None,
                 fonts: dict, log: Callable) -> Path:
    if out_dir is None:
        out_dir = path.parent
    handler = HANDLERS.get(path.suffix.lower())
    if handler is None:
        raise ValueError(f'不支持的格式: {path.suffix}')
    return handler(path, out_dir, fonts, log)


# ── 字体安装 ──────────────────────────────────────────────────────────────────

def ensure_fonts_installed(fonts_dir: Path | None = None) -> list[str]:
    if fonts_dir is None:
        fonts_dir = Path(__file__).parent / 'fonts'
    if not fonts_dir.exists():
        return []
    dst_dir = Path.home() / '.local' / 'share' / 'fonts' / 'office-fix'
    dst_dir.mkdir(parents=True, exist_ok=True)
    installed = []
    for src in fonts_dir.iterdir():
        if src.suffix.lower() not in ('.ttf', '.ttc', '.otf'):
            continue
        dst = dst_dir / src.name
        if not dst.exists():
            shutil.copy2(src, dst)
            installed.append(src.name)
    if installed:
        subprocess.run(['fc-cache', '-f', str(dst_dir)],
                       capture_output=True, check=False)
    return installed
